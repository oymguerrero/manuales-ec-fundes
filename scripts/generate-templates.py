#!/usr/bin/env python3
"""
Genera los 13 templates de los productos del estándar Implementar IA:
- 9 Word (.doc) + 3 Excel (.xls) como HTML estilizado (Office los abre
  nativamente con warning una sola vez por archivo).
- 1 PowerPoint (.pptx) real generado con python-pptx — porque
  PowerPoint moderno NO abre HTML como .ppt, hay que generar OOXML.

Filosofía pedagógica: NO darle la solución al aspirante. Cada template
muestra el criterio F21 oficial + un ejemplo breve del caso La Espiga
(solo para ilustrar el TIPO de contenido) + preguntas guía que el
aspirante responde con datos de SU MiPyME real.

Uso:
    python3 scripts/generate-templates.py

Salida: estandar-a/templates/<num>-<slug>.{doc|xls|ppt}
"""

import os
from pathlib import Path

ROOT = Path(__file__).parent.parent
OUT = ROOT / "estandar-a" / "templates"
OUT.mkdir(parents=True, exist_ok=True)

# Cada producto define su formato óptimo basado en la naturaleza del contenido:
#   "word"  → reportes narrativos, manuales técnicos, actas (HTML en .doc)
#   "excel" → matrices, cronogramas, tablas comparativas (HTML en .xls)
#   "pptx"  → presentación real generada con python-pptx (capacitación)
PRODUCTS = [
    {
        "num": "1.4.1", "slug": "reporte-evaluacion-inicial", "format": "word",
        "elemento": "1 · Planear", "title": "Reporte de evaluación inicial de la MiPyME",
        "intro": "Es el primer producto del proyecto: una foto del estado actual de la MiPyME, sus objetivos y los recursos disponibles.",
        "f21": [
            "Contiene el perfil empresarial de la MiPyME, giro y número de empleados",
            "Incluye el nombre de la MiPyME, ubicación, email y teléfono de contacto",
            "Contiene el resultado de las acciones aplicadas previamente por el cliente",
            "Contiene las áreas de mejora de oportunidades/problemas a resolver desde la perspectiva del cliente y antigüedad de ello",
            "Incluye los objetivos, indicadores y metas",
            "Indica los recursos, tiempos y entregas esperadas",
            "Se presenta de forma digital/físico sin errores ortográficos",
        ],
        "espiga": "Panadería La Espiga (caso pedagógico del manual): 3 sucursales en una ciudad media, 12 empleados, vende pan tradicional. Doña Beatriz (dueña) y Carlos (hijo, apoyo administrativo) quieren reducir el tiempo que las encargadas pasan respondiendo consultas repetitivas por WhatsApp.",
        "preguntas": [
            ("Perfil empresarial", "¿Cómo describirías a la MiPyME en 2-3 líneas? Giro, número de empleados, antigüedad."),
            ("Datos de contacto", "Nombre legal, ubicación, persona responsable, correo y teléfono."),
            ("Acciones previas del cliente", "¿Qué ha intentado hacer la MiPyME antes para resolver estos problemas? ¿Con qué resultados?"),
            ("Áreas de mejora desde la perspectiva del cliente", "¿Qué problemas u oportunidades identifica la MiPyME? ¿Desde cuándo arrastran cada uno?"),
            ("Objetivos · indicadores · metas", "¿Qué quiere lograr la MiPyME con este proyecto? Para cada objetivo, ¿qué indicador lo mide y cuál es la meta?"),
            ("Recursos · tiempos · entregas esperadas", "¿Con qué presupuesto, equipo humano y tiempos cuenta? ¿Qué entregables espera al cierre?"),
        ],
    },
    {
        "num": "1.4.2", "slug": "informe-disponibilidad-datos", "format": "word",
        "elemento": "1 · Planear", "title": "Informe de disponibilidad y calidad de los datos",
        "intro": "Documenta qué datos tiene la MiPyME para sostener soluciones de IA y cuál es su calidad.",
        "f21": [
            "Contiene las fuentes de información existentes en los procesos operativos",
            "Indica que la información requerida está completa y disponible en tiempo",
            "Incluye las necesidades de optimización/estructuración de los datos previos a la implementación",
            "Contiene los hallazgos detectados en la disponibilidad y calidad de los datos",
        ],
        "espiga": "En La Espiga: pedidos por WhatsApp (sin estructura), libretas de papel por sucursal, contabilidad externa en Excel, ningún CRM. Datos dispersos, sin estandarización entre sucursales.",
        "preguntas": [
            ("Fuentes de información existentes", "¿Dónde viven los datos relevantes para el proceso a intervenir? CRM, ERP, hojas de cálculo, mensajes, correos, papel…"),
            ("Completitud y disponibilidad en tiempo", "Por cada fuente: ¿los datos están completos? ¿se actualizan a tiempo? ¿qué frecuencia?"),
            ("Necesidades de optimización/estructuración", "¿Qué hay que limpiar, estandarizar, deduplicar o estructurar antes de poder usarlos con IA?"),
            ("Hallazgos relevantes", "Lo más notable que descubriste sobre la calidad de los datos. ¿Qué obstáculos representa esto para las soluciones que quieres implementar?"),
        ],
    },
    {
        "num": "1.4.3", "slug": "informe-diagnostico-procesos", "format": "word",
        "elemento": "1 · Planear", "title": "Informe de diagnóstico de procesos, actividades y áreas",
        "intro": "Mapea los procesos operativos y detecta los cuellos de botella + resistencias al cambio.",
        "f21": [
            "Contiene el mapa de los procesos analizados",
            "Incluye la clasificación de los procesos por área funcional",
            "Contiene los cuellos de botella y los principales puntos de dolor identificados",
            "Menciona las resistencias al cambio detectadas durante el diagnóstico",
        ],
        "espiga": "Mapa de La Espiga: 6 procesos analizados (recepción de pedidos, producción, venta en mostrador, control de inventario, cobro, contabilidad). Cuello clave: las encargadas pierden 4-5 horas/semana atendiendo consultas repetitivas en WhatsApp. Resistencia: dos panaderos veteranos desconfían de 'que la máquina les diga qué hornear'.",
        "preguntas": [
            ("Mapa de procesos analizados", "Dibuja o describe los procesos que analizaste. Indica qué hace cada uno y cómo se conectan. (Inserta una imagen del mapa si lo dibujaste en Miro/Lucidchart/papel)."),
            ("Clasificación por área funcional", "Agrupa los procesos: ventas, operación, administración, soporte, atención a cliente… ¿Cuál concentra más procesos críticos?"),
            ("Cuellos de botella y puntos de dolor", "Por cada proceso analizado: ¿dónde se atasca? ¿cuánto tiempo se pierde? ¿qué errores ocurren? ¿quién sufre el impacto?"),
            ("Resistencias al cambio detectadas", "Durante el diagnóstico, ¿qué personas mostraron resistencia? ¿Qué temen perder o no entienden? Es información valiosa para la fase de capacitación."),
        ],
    },
    {
        "num": "1.4.4", "slug": "reporte-madurez-digital", "format": "word",
        "elemento": "1 · Planear", "title": "Reporte de madurez digital y disposición al cambio",
        "intro": "Evalúa qué tan preparada está la MiPyME (tecnología) y su personal (cultura) para incorporar IA.",
        "f21": [
            "Incluye la descripción de la metodología y la herramienta utilizadas",
            "Contiene los resultados del test aplicado",
            "Incluye el nivel de madurez digital por dimensión evaluada",
            "Contiene descrito el nivel de disposición al cambio por parte del personal",
            "Incluye recomendaciones basadas en los resultados obtenidos",
            "Incluye un tablero base de indicadores para el seguimiento de la transformación digital",
        ],
        "espiga": "La Espiga obtuvo nivel 2 de 5 en madurez digital (uso de WhatsApp sí, CRM no, datos integrados no). Disposición al cambio: alta en Carlos (28 años), media en las 3 encargadas (40-55), baja en Doña Beatriz al inicio del proyecto (se 'destrabó' tras ver el primer piloto del chatbot).",
        "preguntas": [
            ("Metodología y herramienta usada", "¿Qué instrumento aplicaste para medir madurez digital? (Cuestionario, entrevistas, observación, test estandarizado). ¿Por qué lo elegiste?"),
            ("Resultados del test aplicado", "Resultados crudos por dimensión: infraestructura, procesos, personas, datos, cultura. Apóyate de una tabla."),
            ("Nivel de madurez digital por dimensión", "Clasificación (inicial, en desarrollo, intermedio, avanzado) por cada dimensión. ¿Cuál es la dimensión más débil? ¿Cuál es la fortaleza?"),
            ("Disposición al cambio del personal", "Por persona o grupo: ¿qué tan abierta está al cambio? ¿Qué evidencia observaste durante el diagnóstico?"),
            ("Recomendaciones basadas en resultados", "Acciones específicas para mejorar la madurez en las dimensiones bajas + estrategia para destrabar las resistencias detectadas."),
            ("Tablero base de indicadores", "¿Qué indicadores propondrás seguir mes a mes? Define mínimo 3-5 indicadores con su línea base actual + meta esperada a 6/12 meses."),
        ],
    },
    {
        "num": "1.4.5", "slug": "matriz-impacto-viabilidad", "format": "excel",
        "elemento": "1 · Planear", "title": "Matriz de impacto, viabilidad y esfuerzo",
        "intro": "Hoja de cálculo para priorizar las oportunidades de IA detectadas. Asigna puntajes (1-5) por impacto, viabilidad y esfuerzo; la matriz revela cuáles atacar primero.",
        "f21": [
            "Contiene descritas las soluciones de IA para su aplicación en los procesos priorizados",
            "Incluye las tecnologías/herramientas/plataformas sugeridas de IA adecuadas al contexto de la empresa",
            "Contiene la priorización de las oportunidades de soluciones de IA",
            "Incluye los beneficios esperados para la MiPyME y la viabilidad de la implementación",
        ],
        "espiga": "En La Espiga se identificaron 5 oportunidades. Las dos con mayor impacto + viabilidad (chatbot WhatsApp para FAQ + reporte automatizado de pedidos a cocina) fueron las recomendadas para primera fase.",
        "preguntas": [
            ("Oportunidades detectadas", "Lista todas las oportunidades de IA que detectaste en el diagnóstico. Una por fila."),
            ("Herramienta/plataforma sugerida", "Por cada oportunidad: ¿qué herramienta de IA existente la resolvería? (ChatGPT, Claude, Zapier, Make, CRM con módulo de IA, etc.)"),
            ("Puntaje de IMPACTO (1-5)", "¿Cuánto valor genera? (Tiempo ahorrado, ingresos adicionales, errores reducidos, mejor experiencia). 5 = alto."),
            ("Puntaje de VIABILIDAD (1-5)", "¿Qué tan fácil es de implementar dados los recursos actuales de la MiPyME? 5 = muy viable."),
            ("Puntaje de ESFUERZO (1-5)", "¿Cuánto trabajo cuesta? (Tiempo + personas + costo). 5 = poco esfuerzo, 1 = mucho esfuerzo."),
            ("Beneficios esperados", "Breve frase por oportunidad: ¿qué mejora medible esperas?"),
            ("Priorización final", "Las 2-3 oportunidades recomendadas para primera fase, y por qué."),
        ],
    },
    {
        "num": "1.4.6", "slug": "hoja-de-ruta", "format": "excel",
        "elemento": "1 · Planear", "title": "Hoja de ruta de implementación de soluciones de IA",
        "intro": "Cronograma tipo Gantt con las fases, actividades, responsables y recursos del proyecto.",
        "f21": [
            "Contiene descritas las soluciones de IA seleccionadas para la MiPyME",
            "Contiene descritas la operación de los procesos con la solución de IA",
            "Incluye las fases/etapas de implementación del proyecto",
            "Incluye las actividades y responsables de cada etapa",
            "Contiene el cronograma con las actividades de implementación para cada etapa del proyecto",
            "Incluye los recursos tecnológicos, humanos y financieros requeridos",
            "Incluye indicadores para evaluar el avance de la implementación y darle seguimiento al proyecto",
        ],
        "espiga": "En La Espiga: 4 fases en 5 semanas. S1 configuración del chatbot + plantillas de respuestas FAQ; S2 configuración del reporte automatizado de pedidos; S3 piloto con sucursal principal; S4 expansión a las 3 sucursales + capacitación; S5 evaluación y ajustes.",
        "preguntas": [
            ("Soluciones seleccionadas (de la matriz 1.4.5)", "Lista las soluciones que efectivamente vas a implementar."),
            ("Operación del proceso CON la solución", "Por cada solución: descripción de cómo quedará el proceso después de implementada."),
            ("Fases del proyecto", "¿Cuántas fases tendrá el proyecto? Para cada una: nombre + objetivo + tiempo estimado."),
            ("Actividades y responsables por fase", "Lista granular de actividades. Quién hace qué de tu lado + del lado de la MiPyME."),
            ("Cronograma (Gantt)", "Tabla con actividades en filas y semanas/quincenas en columnas. Marca con X cuándo se ejecuta cada actividad."),
            ("Recursos tecnológicos / humanos / financieros", "Equipo, suscripciones, plataformas, equipo humano, presupuesto desglosado por fase."),
            ("Indicadores de seguimiento", "¿Cómo sabrás semana a semana si vas en tiempo y con calidad? Define indicadores de proceso."),
        ],
    },
    {
        "num": "1.4.7", "slug": "propuesta-final-implementacion", "format": "word",
        "elemento": "1 · Planear", "title": "Propuesta final de implementación integrada",
        "intro": "Documento maestro de entrega al emprendedor que integra los productos previos en una narrativa coherente.",
        "f21": [
            "Contiene el informe de diagnóstico de procesos, actividades y áreas de trabajo (1.4.3)",
            "Incluye el reporte de madurez digital y disposición al cambio (1.4.4)",
            "Contiene la matriz de impacto, viabilidad y esfuerzo (1.4.5)",
            "Incluye la hoja de ruta de implementación de soluciones de IA (1.4.6)",
            "Contiene la propuesta económica de consultoría (1.4.8)",
            "Se presenta en formato digital/físico, con redacción clara y sin errores ortográficos",
        ],
        "espiga": "Para La Espiga: documento maestro de 18-25 páginas. Resumen ejecutivo de 1 pág → diagnóstico → madurez → matriz → hoja de ruta → propuesta económica → anexos. Pensado para que Doña Beatriz lo lea en 20 minutos y entienda exactamente qué pasará.",
        "preguntas": [
            ("Resumen ejecutivo (1 pág)", "Página única que sintetiza: contexto + oportunidades priorizadas + propuesta + inversión + resultados esperados. Si el cliente solo lee esta página, debe entender qué va a pasar."),
            ("Diagnóstico integrado (referencia a 1.4.3)", "Incorpora el informe de diagnóstico o haz un resumen + adjunta el original como anexo."),
            ("Madurez digital y disposición al cambio (referencia a 1.4.4)", "Incluye el reporte o resumen + tablero base."),
            ("Matriz de oportunidades (referencia a 1.4.5)", "Incluye la matriz priorizada — destaca las oportunidades recomendadas para primera fase."),
            ("Hoja de ruta (referencia a 1.4.6)", "Cronograma maestro visible + recursos requeridos."),
            ("Propuesta económica (referencia a 1.4.8)", "Esquema de inversión + propuesta económica detallada — termina llevando al visto bueno por escrito."),
            ("Anexos", "Documentos de soporte: instrumentos del diagnóstico, capturas, evidencias."),
        ],
    },
    {
        "num": "1.4.8", "slug": "propuesta-economica-validada", "format": "word",
        "elemento": "1 · Planear", "title": "Propuesta económica de consultoría validada",
        "intro": "Documento contractual con la inversión + dos REQUISITOS TAXATIVOS del F21: visto bueno por escrito + acuerdo de confidencialidad.",
        "f21": [
            "Contiene la metodología para la implementación de soluciones de IA",
            "Incluye el cronograma, actividades y áreas de trabajo",
            "Contiene descrita la propuesta del monto de inversión por la aplicación del diagnóstico y el test de madurez digital y disposición al cambio",
            "Incluye la evidencia POR ESCRITO del visto bueno por parte de la persona responsable de la MiPyME (REQUISITO TAXATIVO)",
            "Incluye el acuerdo de confidencialidad (REQUISITO TAXATIVO)",
        ],
        "espiga": "Para La Espiga: honorarios profesionales + suscripción mensual de herramientas. Doña Beatriz firmó visto bueno el 15 de marzo + NDA simple de 1 página. Sin esos documentos firmados el evaluador NO da el producto por cumplido.",
        "preguntas": [
            ("Metodología para la implementación", "¿Con qué enfoque vas a ejecutar el proyecto? Pasos generales, principios que guiarán las decisiones."),
            ("Cronograma + áreas de trabajo + actividades", "Síntesis del cronograma + en qué área de la MiPyME ocurre cada actividad."),
            ("Propuesta económica detallada", "Honorarios profesionales (consultor) + costos de licencias/herramientas + otros costos. Total + forma de pago + hitos de cobro si aplica."),
            ("Visto bueno POR ESCRITO del emprendedor", "¿Cómo obtuviste el visto bueno? (Firma autógrafa, electrónica, correo de aprobación explícita, acuse). Adjunta o describe el documento."),
            ("Acuerdo de confidencialidad firmado", "Convenio entre las partes. Adjunta o describe el acuerdo y sus firmas."),
        ],
    },
    {
        "num": "3.4.1", "slug": "soluciones-implementadas", "format": "word",
        "elemento": "2 · Ejecutar", "title": "Inventario de soluciones de IA implementadas",
        "intro": "El producto 3.4.1 NO es documental: es la solución MISMA funcionando. Este template sirve como inventario que documenta lo implementado para la evaluación.",
        "f21": [
            "Se encuentran configuradas en el entorno tecnológico de la MiPyME conforme al plan de implementación definido",
            "Operan en al menos un proceso administrativo/operativo/comercial del negocio de la MiPyME",
            "Incluyen la descripción del caso de uso para el cual fueron implementadas",
            "Contiene evidencia de funcionamiento en condiciones reales de operación",
            "Especifican las herramientas/plataformas de IA utilizadas",
            "Indican el nombre del usuario responsable de la operación dentro de la MiPyME",
        ],
        "espiga": "Inventario para La Espiga: 2 soluciones operando. (1) Chatbot WhatsApp para FAQ desde el 1 de abril, atiende ~80 consultas/día. (2) Reporte automatizado de pedidos del día desde el 5 de abril, ahorra 30 min/día.",
        "preguntas": [
            ("Solución implementada · nombre y proceso intervenido", "Por cada solución: nombre comercial + proceso operativo donde funciona ahora."),
            ("Caso de uso documentado", "¿Para qué se implementó? ¿Qué problema concreto resuelve?"),
            ("Evidencia de funcionamiento", "Capturas, registros de uso, métricas de las últimas 2-4 semanas. (Adjunta como anexos si son muchas)."),
            ("Herramienta/plataforma de IA utilizada", "Nombre comercial exacto + módulo o versión."),
            ("Usuario responsable de la operación", "Quién en la MiPyME opera cada solución (nombre + cargo)."),
            ("Cumplimiento del plan de implementación (1.4.6)", "Cada solución, ¿se desplegó conforme al cronograma? Si hubo desviaciones, ¿por qué?"),
        ],
    },
    {
        "num": "3.4.2", "slug": "informe-tecnico-configuracion", "format": "word",
        "elemento": "2 · Ejecutar", "title": "Informe técnico de configuración",
        "intro": "Manual técnico para que la MiPyME mantenga las soluciones después de tu consultoría. Es exhaustivo.",
        "f21": [
            "Contiene la descripción de las herramientas/plataformas de IA configuradas y su propósito",
            "Incluye los parámetros de configuración aplicados en cada herramienta",
            "Contiene especificadas las integraciones realizadas con sistemas/plataformas/procesos existentes",
            "Incluye los requisitos técnicos necesarios para el funcionamiento",
            "Incluye los elementos de ciberseguridad considerados + leyes y normas aplicables",
            "Contiene los accesos y permisos configurados para los usuarios de la MiPyME",
            "Contiene los resultados de las pruebas funcionales (funcionamiento/integración/uso en condiciones reales)",
            "Contiene especificados los indicadores de desempeño obtenidos durante la fase de pruebas",
            "Se presenta en formato digital, con redacción clara y sin errores ortográficos",
        ],
        "espiga": "Informe técnico de La Espiga: 14 páginas que documentan cada parámetro del chatbot WhatsApp, la integración con Google Sheets para registrar pedidos, los permisos de las 3 encargadas + Carlos, y los resultados de 15 pruebas funcionales ejecutadas antes del go-live.",
        "preguntas": [
            ("Resumen ejecutivo", "Una página: qué se implementó, en qué entorno, con qué arquitectura general."),
            ("Arquitectura general", "Diagrama de cómo se conectan las herramientas entre sí y con los procesos de la MiPyME."),
            ("Herramientas + propósito + parámetros", "Por cada herramienta: nombre comercial, propósito en la MiPyME, todos los parámetros de configuración aplicados."),
            ("Integraciones realizadas", "¿Con qué sistemas existentes se conectó cada herramienta? (CRM, hojas de cálculo, plataformas de mensajería, etc.)"),
            ("Requisitos técnicos", "Conectividad, navegadores, sistemas operativos, suscripciones que la MiPyME necesita mantener."),
            ("Ciberseguridad + normativa aplicable", "Medidas de seguridad implementadas + cumplimiento LFPDPPP (mínimo Arts. 5-12 principios + 7,8,14 aviso de privacidad + 36-37 transferencias + 22-34 ARCO + 63-68 sanciones)."),
            ("Accesos y permisos configurados", "Tabla con cada usuario + rol + nivel de acceso por herramienta."),
            ("Pruebas funcionales ejecutadas", "Lista de pruebas (funcionamiento, integración, uso en condiciones reales) + resultados de cada una."),
            ("Indicadores de desempeño en pruebas", "Métricas obtenidas durante las pruebas: tiempo de respuesta, precisión, errores, etc."),
        ],
    },
    {
        "num": "3.4.3", "slug": "material-capacitacion", "format": "pptx",
        "elemento": "2 · Ejecutar", "title": "Material de capacitación al personal",
        "intro": "Presentación visual para capacitar al personal de la MiPyME. Una lámina por paso, espacio para captura, prompts destacados y errores frecuentes. El formato PowerPoint te permite usarla en sesión presencial o compartirla para autoestudio.",
        "f21": [
            "Contiene las instrucciones paso a paso para operar cada herramienta de IA implementada",
            "Incluye ejemplos de uso aplicados al contexto específico de la MiPyME",
            "Especifica los prompts o instrucciones recomendadas para cada herramienta de IA",
            "Incluye las buenas prácticas de uso y las acciones a realizar ante errores frecuentes",
            "Considera el nivel de alfabetización digital del personal capacitado",
            "Se presenta en formato digital, accesible y comprensible para el personal de la MiPyME",
        ],
        "espiga": "Material de La Espiga: presentación de 12 láminas + 3 prompts pre-cargados. Pensada para encargadas con uso intermedio de smartphone pero sin experiencia previa con asistentes de IA. Lenguaje sencillo, capturas grandes, lista de 'qué hacer si algo no funciona'.",
        "preguntas": [
            ("Lámina 1 · Bienvenida y objetivos del material", "¿Qué van a aprender? ¿Cuánto tiempo tomará? ¿Para qué les servirá en su día a día?"),
            ("Láminas 2-N · Instrucciones paso a paso por herramienta", "Por cada herramienta: capturas numeradas con el flujo completo de uso. Si la herramienta tiene 5 botones importantes, dedica 1 lámina por botón. Visual antes que texto."),
            ("Ejemplos aplicados al contexto de la MiPyME", "Por cada herramienta: 2-3 ejemplos REALES de la MiPyME (no ejemplos genéricos). El personal debe poder decir 'ah, eso lo hago yo el lunes en la mañana'."),
            ("Prompts recomendados", "Lista de 5-10 prompts que ya funcionan para los casos típicos. El personal los copia/pega; no se les pide redactar prompts desde cero."),
            ("Errores frecuentes y qué hacer", "Tabla con: error común → cómo identificarlo → cómo resolverlo → cuándo escalar al consultor o a Carlos/Doña Beatriz."),
            ("Adaptación al nivel digital del personal", "Si el personal no maneja tecnicismos, usa lenguaje cotidiano. Si están acostumbrados a apps, puedes ir más rápido."),
            ("Lámina final · Recursos y soporte", "¿A quién acudir si algo falla? ¿Dónde encontrar este material después? ¿Cuándo se actualiza?"),
        ],
    },
    {
        "num": "4.4.1", "slug": "reporte-evaluacion-resultados", "format": "excel",
        "elemento": "3 · Evaluar", "title": "Reporte de evaluación de resultados",
        "intro": "Hoja de cálculo con la comparación antes/después de cada indicador + tablero visual + estimación de ROI.",
        "f21": [
            "Contiene la comparación de los indicadores de desempeño de los procesos antes y después de la implementación",
            "Incluye el tablero de indicadores integrando resultados cuantitativos y cualitativos en los procesos intervenidos",
            "Contiene el impacto operativo de las soluciones (tiempos, Retorno de inversión, errores reducidos)",
            "Incluye el registro de incidencias presentadas durante la operación y las acciones correctivas aplicadas",
            "Contiene especificadas el nivel de adopción de las soluciones por parte del personal de la MiPyME",
            "Se presenta en formato digital/físico, con redacción clara y sin errores ortográficos",
        ],
        "espiga": "Reporte de La Espiga: tabla comparativa de 5 indicadores (tiempo de respuesta WhatsApp, consultas atendidas/día, horas semanales en tareas repetitivas, errores en pedidos, ventas online). Resultado: ROI estimado 280% a 6 meses. 1 incidencia mayor (caída del chatbot 2 horas el 18 de abril) con su acción correctiva documentada.",
        "preguntas": [
            ("Indicadores · línea base · valor actual · variación", "Tabla con cada indicador del proyecto, su valor antes de la implementación, su valor actual, y la variación porcentual."),
            ("Tablero de indicadores integrado", "Visualización gráfica (puedes generar gráficas con Excel mismo) con resultados cuantitativos + cualitativos."),
            ("Impacto operativo", "Tiempo ahorrado, ingresos adicionales, errores reducidos. Cuantifica todo lo cuantificable."),
            ("Cálculo del ROI", "Beneficios obtenidos / costo total del proyecto × 100. Documenta los supuestos."),
            ("Registro de incidencias + acciones correctivas", "Tabla con cada incidencia: fecha, qué ocurrió, qué hiciste para corregirla, qué se aprendió."),
            ("Nivel de adopción del personal", "¿Cuántas personas operan cada solución activamente? ¿Qué porcentaje del personal capacitado la usa? Evidencia: logs, encuesta a usuarios."),
            ("Lecciones y recomendaciones", "Resumen ejecutivo de lo que funcionó, lo que no, y recomendaciones de mejora continua."),
        ],
    },
    {
        "num": "4.4.2", "slug": "acta-de-cierre", "format": "word",
        "elemento": "3 · Evaluar", "title": "Acta de cierre del proyecto",
        "intro": "Documento legal de cierre. CRÍTICO: requiere firmas de ambas partes + contactos de soporte. Sin estos elementos el evaluador NO da por cumplido el producto.",
        "f21": [
            "Contiene el resumen ejecutivo del proyecto: alcance, objetivos y resultados obtenidos",
            "Incluye el listado de las soluciones de IA implementadas y su estado operativo al cierre",
            "Incluye las lecciones aprendidas durante el proyecto de implementación",
            "Contiene descritos los compromisos pendientes y próximos pasos acordados entre ambas partes",
            "Contiene los contactos de soporte técnico para incidencias futuras (REQUISITO TAXATIVO)",
            "Incluye las firmas del consultor y del representante de la MiPyME (REQUISITO TAXATIVO)",
            "Se presenta en formato digital/físico, sin errores ortográficos",
        ],
        "espiga": "Acta de La Espiga firmada el 30 de mayo por Doña Beatriz y por el consultor. Compromisos pendientes: revisar el chatbot en julio para añadir 10 nuevas FAQ. Contacto de soporte: correo + WhatsApp del consultor. Lección clave registrada: 'capacitar PRIMERO a las personas más abiertas al cambio, luego ellas convencen a las escépticas'.",
        "preguntas": [
            ("Resumen ejecutivo del proyecto", "Alcance original, objetivos planteados, resultados obtenidos. Una página."),
            ("Soluciones implementadas + estado al cierre", "Lista de cada solución + estado: operando / operando con incidencias / pausada / cancelada."),
            ("Lecciones aprendidas", "Lo que funcionó, lo que no, lo que harías diferente. 5-7 lecciones específicas."),
            ("Compromisos pendientes y próximos pasos", "Si quedó algo a futuro: qué, cuándo, responsable. Si no quedó nada pendiente, decláralo explícitamente."),
            ("Contactos de soporte técnico (REQUISITO TAXATIVO)", "Datos completos del consultor para incidencias futuras: nombre, correo, teléfono, horario de respuesta esperado."),
            ("Firmas (REQUISITO TAXATIVO)", "Firma del consultor + firma del representante legal o emprendedor de la MiPyME. Físicas o digitales con valor legal."),
        ],
    },
]


# ===========================================================================
# PLANTILLAS POR FORMATO
# ===========================================================================

BASE_STYLES = """
@page {
  size: 21.59cm 27.94cm;
  margin: 2.2cm 1.8cm;
}
body {
  font-family: 'Calibri', 'Arial', sans-serif;
  font-size: 11pt;
  line-height: 1.5;
  color: #333;
}
.brand-header {
  border-bottom: 3pt solid #f7c031;
  padding-bottom: 10pt;
  margin-bottom: 18pt;
}
.brand-header table { width: 100%; border-collapse: collapse; }
.brand-header td { border: none; vertical-align: middle; padding: 0; }
.brand-block {
  background: #28467e;
  color: #f7c031;
  font-weight: bold;
  font-size: 13pt;
  padding: 12pt 8pt;
  text-align: center;
  border-radius: 6pt;
  width: 70pt;
}
.brand-header h1 { color: #28467e; font-size: 17pt; margin: 0 0 3pt 0; font-weight: bold; }
.brand-header .subtitle { color: #666; font-size: 10pt; margin: 0; }
.brand-header .meta { text-align: right; font-size: 9pt; color: #666; width: 130pt; }
.brand-header .meta strong { color: #28467e; font-size: 11pt; }

.product-data {
  background: #f5f9ff;
  border: 1pt solid #cce0f5;
  padding: 12pt 16pt;
  margin: 0 0 20pt 0;
  border-radius: 4pt;
}
.product-data table { width: 100%; border-collapse: collapse; }
.product-data td { padding: 4pt 8pt; font-size: 10pt; border: none; }
.product-data td:first-child { font-weight: bold; color: #28467e; width: 30%; }
.product-data .blank { border-bottom: 1pt solid #999; display: inline-block; min-width: 220pt; }

.intro {
  background: #FFF3CC;
  border-left: 4pt solid #f7c031;
  padding: 12pt 16pt;
  margin: 14pt 0 20pt;
  font-size: 10pt;
  color: #555;
}

.f21-box {
  background: #f0f7ff;
  border-left: 4pt solid #1F4E8C;
  padding: 12pt 16pt 14pt;
  margin: 16pt 0;
  font-size: 10pt;
}
.f21-box h3 {
  margin: 0 0 8pt 0;
  font-size: 11pt;
  color: #28467e;
  text-transform: uppercase;
  letter-spacing: 0.04em;
}
.f21-box ul { margin: 0; padding-left: 18pt; color: #333; }
.f21-box li { margin-bottom: 4pt; }
.f21-box .taxativo { background: #FFE0E0; color: #C0392B; font-weight: bold; padding: 1pt 6pt; border-radius: 3pt; font-size: 9pt; margin-left: 4pt; }

.espiga-box {
  background: #fdf6e3;
  border: 1pt dashed #d4a40e;
  padding: 12pt 16pt;
  margin: 14pt 0 20pt;
  font-size: 10pt;
  font-style: italic;
  color: #6a5208;
}
.espiga-box strong { color: #28467e; font-style: normal; }

h2 {
  color: #28467e;
  font-size: 14pt;
  border-bottom: 1pt solid #f7c031;
  padding-bottom: 4pt;
  margin-top: 24pt;
  margin-bottom: 12pt;
}

.question-block {
  margin-bottom: 22pt;
  padding: 0 0 12pt 0;
  border-bottom: 1pt dotted #ddd;
}
.question-num {
  display: inline-block;
  background: #28467e;
  color: white;
  width: 22pt;
  height: 22pt;
  line-height: 22pt;
  text-align: center;
  border-radius: 50%;
  font-weight: bold;
  font-size: 11pt;
  margin-right: 8pt;
  vertical-align: middle;
}
.question-title {
  font-weight: bold;
  color: #28467e;
  font-size: 11.5pt;
}
.question-help {
  color: #666;
  font-size: 10pt;
  font-style: italic;
  margin: 6pt 0 8pt 30pt;
  line-height: 1.45;
}
.fill-area {
  border: 1pt dashed #ccc;
  background: #fafafa;
  padding: 22pt 16pt;
  margin: 6pt 0 0 30pt;
  min-height: 60pt;
  color: #aaa;
  font-style: italic;
  font-size: 10pt;
}

.footer {
  margin-top: 36pt;
  padding-top: 10pt;
  border-top: 1pt solid #ddd;
  font-size: 8pt;
  color: #999;
  text-align: center;
  font-style: italic;
}
"""

WORD_TEMPLATE = """<!DOCTYPE html>
<html xmlns:o="urn:schemas-microsoft-com:office:office"
      xmlns:w="urn:schemas-microsoft-com:office:word"
      xmlns="http://www.w3.org/TR/REC-html40">
<head>
<meta http-equiv="Content-Type" content="text/html; charset=utf-8">
<meta name="ProgId" content="Word.Document">
<meta name="Generator" content="Microsoft Word">
<title>{num} · {title}</title>
<!--[if gte mso 9]>
<xml><w:WordDocument><w:View>Print</w:View><w:Zoom>100</w:Zoom><w:DoNotOptimizeForBrowser/></w:WordDocument></xml>
<![endif]-->
<style>{base_styles}</style>
</head>
<body>

<div class="brand-header">
<table>
<tr>
<td style="width: 80pt;"><div class="brand-block">Mi<br>CompañIA</div></td>
<td style="padding-left: 14pt;">
<h1>{title}</h1>
<p class="subtitle">Manual de Implementar IA · Producto del Elemento {elemento}</p>
</td>
<td class="meta"><strong>Producto {num}</strong><br>Template editable<br><em>Mi CompañIA · FUNDES</em></td>
</tr>
</table>
</div>

<div class="intro"><strong>Sobre este template.</strong> {intro} Tu evaluador no busca que llenes el template — busca que respondas con datos de TU MiPyME real. Las preguntas guía te orientan; las respuestas son tuyas.</div>

<div class="f21-box">
<h3>Qué evalúa el F21 oficial</h3>
<ul>{f21_html}</ul>
</div>

<div class="espiga-box"><strong>Caso pedagógico La Espiga (ilustrativo, NO es tu solución):</strong> {espiga}</div>

<div class="product-data">
<table>
<tr><td>MiPyME (razón social)</td><td><span class="blank">&nbsp;</span></td></tr>
<tr><td>Persona responsable</td><td><span class="blank">&nbsp;</span></td></tr>
<tr><td>Consultor responsable</td><td><span class="blank">&nbsp;</span></td></tr>
<tr><td>Fecha del documento</td><td><span class="blank">&nbsp;</span></td></tr>
</table>
</div>

<h2>Tu turno · preguntas guía</h2>

{questions_html}

<div class="footer">
Mi CompañIA · Una iniciativa de FUNDES Latinoamérica con el apoyo de Google.org<br>
Template del Manual de Implementar IA · Para ser llenado por el aspirante con datos de su proyecto real con una MiPyME
</div>

</body>
</html>
"""

EXCEL_TEMPLATE = """<!DOCTYPE html>
<html xmlns:o="urn:schemas-microsoft-com:office:office"
      xmlns:x="urn:schemas-microsoft-com:office:excel"
      xmlns="http://www.w3.org/TR/REC-html40">
<head>
<meta http-equiv="Content-Type" content="text/html; charset=utf-8">
<meta name="ProgId" content="Excel.Sheet">
<meta name="Generator" content="Microsoft Excel">
<title>{num} · {title}</title>
<!--[if gte mso 9]>
<xml><x:ExcelWorkbook><x:ExcelWorksheets>
<x:ExcelWorksheet><x:Name>{num}</x:Name><x:WorksheetOptions><x:Panes><x:Pane><x:Number>3</x:Number></x:Pane></x:Panes></x:WorksheetOptions></x:ExcelWorksheet>
</x:ExcelWorksheets></x:ExcelWorkbook></xml>
<![endif]-->
<style>{base_styles}
table.matrix {{ border-collapse: collapse; width: 100%; margin: 10pt 0; font-size: 10pt; }}
table.matrix th, table.matrix td {{ border: 1pt solid #999; padding: 6pt 8pt; vertical-align: top; }}
table.matrix th {{ background: #28467e; color: white; font-weight: bold; text-align: left; }}
table.matrix tr.score th {{ background: #f7c031; color: #28467e; }}
table.matrix td.fill {{ background: #fffbe6; min-height: 22pt; height: 22pt; color: #aaa; font-style: italic; }}
table.matrix tr:nth-child(even) td.fill {{ background: #fefbf0; }}
</style>
</head>
<body>

<div class="brand-header">
<table>
<tr>
<td style="width: 80pt;"><div class="brand-block">Mi<br>CompañIA</div></td>
<td style="padding-left: 14pt;">
<h1>{title}</h1>
<p class="subtitle">Manual de Implementar IA · Producto del Elemento {elemento} · Excel</p>
</td>
<td class="meta"><strong>Producto {num}</strong><br>Hoja editable · Excel<br><em>Mi CompañIA · FUNDES</em></td>
</tr>
</table>
</div>

<div class="intro"><strong>Sobre este template.</strong> {intro} Excel es el formato natural para esta clase de producto: te permite añadir filas conforme detectes más oportunidades, ordenar por puntaje y generar gráficas. <strong>Llénalo con datos de TU MiPyME real.</strong></div>

<div class="f21-box">
<h3>Qué evalúa el F21 oficial</h3>
<ul>{f21_html}</ul>
</div>

<div class="espiga-box"><strong>Caso pedagógico La Espiga (ilustrativo):</strong> {espiga}</div>

<div class="product-data">
<table>
<tr><td>MiPyME</td><td><span class="blank">&nbsp;</span></td><td>Fecha</td><td><span class="blank">&nbsp;</span></td></tr>
<tr><td>Consultor</td><td><span class="blank">&nbsp;</span></td><td>Responsable MiPyME</td><td><span class="blank">&nbsp;</span></td></tr>
</table>
</div>

<h2>Tu turno · preguntas guía</h2>

{questions_html}

{matrix_html}

<div class="footer">
Mi CompañIA · FUNDES Latinoamérica con el apoyo de Google.org · Template Excel del Manual de Implementar IA
</div>

</body>
</html>
"""

# ===========================================================================
# GENERADOR PPTX REAL (3.4.3 material de capacitación)
# ===========================================================================
#
# python-pptx solo se importa si hay productos con format="pptx".
# La estética: 16:9, marca Mi CompañIA arriba, layout consistente con los
# otros 12 templates HTML (mismo azul, mismo amarillo, mismas cajas).

def generate_pptx_capacitacion(prod):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    AZUL = RGBColor(0x28, 0x46, 0x7E)
    AMARILLO = RGBColor(0xF7, 0xC0, 0x31)
    AZUL_CLARO = RGBColor(0xF0, 0xF7, 0xFF)
    AMARILLO_CLARO = RGBColor(0xFF, 0xF3, 0xCC)
    BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
    NEGRO = RGBColor(0x33, 0x33, 0x33)
    GRIS = RGBColor(0x66, 0x66, 0x66)
    GRIS_CLARO = RGBColor(0xFA, 0xFA, 0xFA)
    GRIS_PUNTEADO = RGBColor(0xCC, 0xE0, 0xF5)
    CREMA = RGBColor(0xFD, 0xF6, 0xE3)
    ROJO = RGBColor(0xC0, 0x39, 0x2B)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def solid(shape, color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def textbox(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
                color=AZUL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font_name
        return tb

    def filled_shape(slide, shape_type, x, y, w, h, fill, line_color=None):
        s = slide.shapes.add_shape(shape_type, x, y, w, h)
        solid(s, fill)
        if line_color is not None:
            s.line.color.rgb = line_color
            s.line.width = Pt(0.75)
        return s

    def header_bar(slide, prod):
        # franja azul superior
        band = filled_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.5), AZUL)
        # franja amarilla justo debajo
        filled_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(0.5), SW, Inches(0.06), AMARILLO)
        # logo "Mi CompañIA"
        textbox(slide, Inches(0.35), Inches(0.05), Inches(3.5), Inches(0.4),
                "Mi CompañIA", size=14, bold=True, color=AMARILLO, anchor=MSO_ANCHOR.MIDDLE)
        # título de producto
        textbox(slide, Inches(3.9), Inches(0.05), Inches(7.5), Inches(0.4),
                f"· {prod['title']}", size=12, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
        # número de producto a la derecha
        textbox(slide, Inches(11.4), Inches(0.05), Inches(1.8), Inches(0.4),
                f"Producto {prod['num']}", size=12, bold=True, color=AMARILLO,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    def footer_band(slide, page_num=None, total=None):
        filled_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(7.15), SW, Inches(0.35), AMARILLO)
        label = "Mi CompañIA  ·  FUNDES Latinoamérica  ·  AIxMiPyMEs"
        if page_num is not None and total is not None:
            label += f"   ·   {page_num} / {total}"
        textbox(slide, Inches(0), Inches(7.2), SW, Inches(0.3),
                label, size=10, bold=True, color=AZUL, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    # cuántas slides totales calcularemos: portada + datos + F21 + Espiga + N preguntas + cierre
    total = 4 + len(prod["preguntas"]) + 1
    page = 1

    # ---- Slide 1: portada ----
    s = prs.slides.add_slide(blank)
    bg = filled_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, AZUL)
    # bloque decorativo amarillo a la derecha
    filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(9.5), 0, Inches(3.833), SH, AMARILLO)
    filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(9.5), 0, Inches(0.12), SH, AMARILLO)
    # cuadro central
    textbox(s, Inches(0.7), Inches(2.0), Inches(8.5), Inches(0.6),
            f"PRODUCTO {prod['num']}  ·  ELEMENTO {prod['elemento']}",
            size=15, bold=True, color=AMARILLO)
    textbox(s, Inches(0.7), Inches(2.7), Inches(8.5), Inches(2.0),
            prod["title"], size=44, bold=True, color=BLANCO)
    # subtitle
    textbox(s, Inches(0.7), Inches(5.0), Inches(8.5), Inches(0.6),
            "Material de capacitación al personal", size=20, italic=True, color=AZUL_CLARO)
    # banda amarilla con autoría
    textbox(s, Inches(0.7), Inches(5.9), Inches(8.5), Inches(0.5),
            "Mi CompañIA  ·  FUNDES Latinoamérica con el apoyo de Google.org",
            size=12, bold=True, color=AMARILLO)
    # logo en banda derecha
    badge = filled_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(3.0),
                        Inches(2.7), Inches(1.5), AZUL)
    tf = badge.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Mi\nCompañIA"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = AMARILLO

    # ---- Slide 2: datos de la capacitación ----
    page += 1
    s = prs.slides.add_slide(blank)
    header_bar(s, prod)
    textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.8),
            "Datos de esta capacitación", size=32, bold=True, color=AZUL)
    textbox(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.6),
            "Completa antes de impartirla. Define alcance, audiencia y tiempo para que el evaluador vea contexto.",
            size=13, italic=True, color=GRIS)
    fields = [
        "Nombre de la MiPyME",
        "Personal capacitado (nombres + cargo)",
        "Fecha de la capacitación",
        "Duración estimada",
        "Soluciones cubiertas en esta sesión",
    ]
    for i, label in enumerate(fields):
        y = Inches(2.7 + i * 0.75)
        # franja azul claro
        filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(0.6), AZUL_CLARO)
        # label
        textbox(s, Inches(0.7), y, Inches(4.0), Inches(0.6),
                label + ":", size=13, bold=True, color=AZUL, anchor=MSO_ANCHOR.MIDDLE)
        # línea para llenar
        textbox(s, Inches(4.8), y, Inches(8.0), Inches(0.6),
                "________________________________________", size=13, color=GRIS, anchor=MSO_ANCHOR.MIDDLE)
    footer_band(s, page, total)

    # ---- Slide 3: criterio F21 ----
    page += 1
    s = prs.slides.add_slide(blank)
    header_bar(s, prod)
    textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.8),
            "Criterio F21 · qué evalúa el verificador",
            size=28, bold=True, color=AZUL)
    textbox(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.5),
            "Texto literal del F21 oficial. Tu material debe cubrir explícitamente cada punto.",
            size=12, italic=True, color=GRIS)
    box = filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.4),
                      Inches(12.3), Inches(4.5), AZUL_CLARO)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)
    tf.margin_right = Inches(0.4)
    for i, crit in enumerate(prod["f21"]):
        is_tax = "REQUISITO TAXATIVO" in crit
        clean = crit.replace("(REQUISITO TAXATIVO)", "").strip()
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = "•  " + clean
        r.font.size = Pt(14)
        r.font.color.rgb = AZUL if not is_tax else ROJO
        r.font.bold = is_tax
        if is_tax:
            r2 = p.add_run()
            r2.text = "   ⚠ REQUISITO TAXATIVO"
            r2.font.size = Pt(11)
            r2.font.bold = True
            r2.font.color.rgb = ROJO
    footer_band(s, page, total)

    # ---- Slide 4: caso La Espiga ----
    page += 1
    s = prs.slides.add_slide(blank)
    header_bar(s, prod)
    textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.8),
            "Ejemplo orientativo · caso La Espiga", size=28, bold=True, color=AZUL)
    textbox(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.5),
            "Solo para mostrar el TIPO de contenido. NO copies esto a tu MiPyME.",
            size=12, italic=True, color=GRIS)
    box = filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.4),
                      Inches(12.3), Inches(4.5), CREMA)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = prod["espiga"]
    r.font.size = Pt(16)
    r.font.italic = True
    r.font.color.rgb = RGBColor(0x6A, 0x52, 0x08)
    footer_band(s, page, total)

    # ---- Slide por cada pregunta guía ----
    for idx, (title, help_text) in enumerate(prod["preguntas"], start=1):
        page += 1
        title_lower = title.lower()
        is_prompt = "prompt" in title_lower
        is_errors = "error" in title_lower or "qué hacer" in title_lower

        s = prs.slides.add_slide(blank)
        header_bar(s, prod)
        # número grande en círculo
        circle = filled_shape(s, MSO_SHAPE.OVAL, Inches(0.55), Inches(0.95),
                             Inches(1.1), Inches(1.1), AZUL)
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(idx)
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = AMARILLO
        # eyebrow "LÁMINA N"
        textbox(s, Inches(1.85), Inches(0.95), Inches(11), Inches(0.4),
                f"LÁMINA {idx}", size=11, bold=True, color=AMARILLO)
        # título
        textbox(s, Inches(1.85), Inches(1.3), Inches(11.2), Inches(0.9),
                title, size=26, bold=True, color=AZUL)
        # subrayado amarillo
        filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(1.85), Inches(2.18),
                    Inches(2.5), Inches(0.06), AMARILLO)
        # help text
        textbox(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.8),
                help_text, size=13, italic=True, color=GRIS)

        # Contenido específico
        if is_prompt:
            # 2 cajas de prompt apiladas
            for i, (label_text, bg_color, accent) in enumerate([
                ("PROMPT SUGERIDO", AMARILLO_CLARO, AMARILLO),
                ("OTRO PROMPT", AZUL_CLARO, AZUL),
            ]):
                y = Inches(3.5 + i * 1.7)
                # franja izquierda como acento
                filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(1.5), accent)
                box = filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(0.62), y, Inches(12.2), Inches(1.5), bg_color)
                tf = box.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.3)
                tf.margin_top = Inches(0.15)
                tf.margin_right = Inches(0.3)
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = label_text
                r.font.size = Pt(10)
                r.font.bold = True
                r.font.color.rgb = AZUL
                p2 = tf.add_paragraph()
                p2.space_before = Pt(4)
                r = p2.add_run()
                r.text = "[ Pega aquí un prompt completo: rol del asistente + contexto de la MiPyME + formato de salida esperado + ejemplo. ]"
                r.font.size = Pt(13)
                r.font.italic = True
                r.font.color.rgb = RGBColor(0x5A, 0x44, 0x00)
                r.font.name = "Consolas"
        elif is_errors:
            # tabla de 4 columnas x 5 filas (1 header + 4 data)
            tbl_shape = s.shapes.add_table(rows=5, cols=4,
                                            left=Inches(0.5), top=Inches(3.5),
                                            width=Inches(12.3), height=Inches(3.5))
            tbl = tbl_shape.table
            tbl.columns[0].width = Inches(2.6)
            tbl.columns[1].width = Inches(2.5)
            tbl.columns[2].width = Inches(3.7)
            tbl.columns[3].width = Inches(3.5)
            headers = ["Error frecuente", "Cómo identificarlo", "Cómo resolverlo", "Cuándo escalar"]
            for c, h in enumerate(headers):
                cell = tbl.cell(0, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = AZUL
                cell.text = h
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(12)
                        run.font.bold = True
                        run.font.color.rgb = BLANCO
            for row in range(1, 5):
                for col in range(4):
                    cell = tbl.cell(row, col)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GRIS_CLARO if row % 2 == 1 else BLANCO
                    cell.text = ""
        else:
            # caja de captura + caja de notas
            cap = filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.5),
                              Inches(7.0), Inches(3.5), AZUL_CLARO, line_color=GRIS_PUNTEADO)
            tf = cap.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.3)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = "📷  Inserta aquí la captura\n\n[Sugerencia: marca con círculo rojo cada zona importante]"
            r.font.size = Pt(13)
            r.font.italic = True
            r.font.color.rgb = RGBColor(0x6D, 0x8D, 0xB6)
            # caja de notas a la derecha
            notes = filled_shape(s, MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(3.5),
                                Inches(5.1), Inches(3.5), GRIS_CLARO)
            tf = notes.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.3)
            tf.margin_top = Inches(0.2)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "NOTAS PARA EL CAPACITADOR"
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = AZUL
            p2 = tf.add_paragraph()
            p2.space_before = Pt(8)
            r = p2.add_run()
            r.text = "[ Escribe aquí qué decir mientras muestras la captura. Ejemplos reales de la MiPyME van mejor que descripciones genéricas. ]"
            r.font.size = Pt(12)
            r.font.italic = True
            r.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        footer_band(s, page, total)

    # ---- Slide final: cierre ----
    page += 1
    s = prs.slides.add_slide(blank)
    filled_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, AZUL)
    textbox(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.5),
            "¿Dudas durante la operación?", size=42, bold=True, color=AMARILLO,
            align=PP_ALIGN.CENTER)
    textbox(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.8),
            "Tu consultor sigue disponible. Datos de contacto en el Acta de cierre (producto 4.4.2).",
            size=18, color=BLANCO, align=PP_ALIGN.CENTER)
    textbox(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
            "Este material puede actualizarse cuando se añadan nuevas FAQ o flujos.",
            size=14, italic=True, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
    filled_shape(s, MSO_SHAPE.RECTANGLE, 0, Inches(6.85), SW, Inches(0.65), AMARILLO)
    textbox(s, Inches(0), Inches(6.95), SW, Inches(0.5),
            "Mi CompañIA  ·  FUNDES Latinoamérica con el apoyo de Google.org",
            size=13, bold=True, color=AZUL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    return prs


# ===========================================================================
# RENDERIZADORES
# ===========================================================================

def render_f21(criterios):
    html = ""
    for c in criterios:
        is_tax = "REQUISITO TAXATIVO" in c
        text = c.replace(" (REQUISITO TAXATIVO)", "")
        if is_tax:
            html += f'<li>{text} <span class="taxativo">REQUISITO TAXATIVO</span></li>\n'
        else:
            html += f'<li>{text}</li>\n'
    return html


def render_questions_word(preguntas):
    html = ""
    for i, (title, help_text) in enumerate(preguntas, 1):
        html += (
            f'<div class="question-block">\n'
            f'<div><span class="question-num">{i}</span><span class="question-title">{title}</span></div>\n'
            f'<div class="question-help">{help_text}</div>\n'
            f'<div class="fill-area">[Escribe aquí tu respuesta con datos reales de tu MiPyME. Borra esta línea cuando empieces.]</div>\n'
            f'</div>\n'
        )
    return html


def render_questions_excel(preguntas):
    html = '<table class="matrix">\n<tr><th style="width: 4%;">#</th><th style="width: 28%;">Pregunta guía</th><th>Tu respuesta (llena con datos de TU MiPyME)</th></tr>\n'
    for i, (title, help_text) in enumerate(preguntas, 1):
        html += (
            f'<tr>\n'
            f'<td style="text-align: center; font-weight: bold; color: #28467e;">{i}</td>\n'
            f'<td><strong>{title}</strong><br><span style="color: #666; font-size: 9pt; font-style: italic;">{help_text}</span></td>\n'
            f'<td class="fill">&nbsp;</td>\n'
            f'</tr>\n'
        )
    html += '</table>\n'
    return html


def render_matrix_example(num):
    """Para 1.4.5 (matriz impacto/viabilidad), añade una hoja de cálculo pre-armada."""
    if num != "1.4.5":
        return ""
    rows = ""
    for i in range(1, 9):
        rows += (
            f'<tr>\n'
            f'<td style="text-align:center;color:#28467e;font-weight:bold;">{i}</td>\n'
            f'<td class="fill">&nbsp;</td><td class="fill">&nbsp;</td>\n'
            f'<td class="fill" style="text-align:center;">&nbsp;</td>\n'
            f'<td class="fill" style="text-align:center;">&nbsp;</td>\n'
            f'<td class="fill" style="text-align:center;">&nbsp;</td>\n'
            f'<td class="fill" style="text-align:center;color:#999;">=I*V/E</td>\n'
            f'<td class="fill">&nbsp;</td>\n'
            f'</tr>\n'
        )
    return f"""
<h2>Hoja de matriz · oportunidades de IA priorizadas</h2>
<p style="font-size: 10pt; color: #666;">Añade tantas filas como oportunidades hayas detectado. Llena impacto (I), viabilidad (V) y esfuerzo (E) con valores 1-5. El score relativo (I × V / E) te orienta sobre cuál priorizar.</p>
<table class="matrix">
<tr><th style="width:4%;">#</th><th style="width:22%;">Oportunidad</th><th style="width:16%;">Herramienta sugerida</th><th style="width:9%;">Impacto (1-5)</th><th style="width:9%;">Viabilidad (1-5)</th><th style="width:9%;">Esfuerzo (1-5)</th><th style="width:9%;">Score</th><th>Beneficio esperado</th></tr>
{rows}
</table>
<p style="font-size:9pt;color:#999;margin-top:8pt;font-style:italic;">Convención: I y V altos suman, esfuerzo alto resta (por eso aparece en el denominador). No hay fórmula universal — adapta el cálculo a tu criterio profesional.</p>
"""


def render_gantt_example(num):
    """Para 1.4.6 (hoja de ruta), añade un cronograma Gantt vacío."""
    if num != "1.4.6":
        return ""
    weeks_header = ''.join(f'<th style="width:5%;text-align:center;">S{i}</th>' for i in range(1, 13))
    rows = ""
    for i in range(1, 9):
        empty_cells = ''.join('<td class="fill" style="text-align:center;">&nbsp;</td>' for _ in range(12))
        rows += (
            f'<tr>\n'
            f'<td style="text-align:center;color:#28467e;font-weight:bold;">{i}</td>\n'
            f'<td class="fill">&nbsp;</td>\n'
            f'<td class="fill">&nbsp;</td>\n'
            f'{empty_cells}\n'
            f'</tr>\n'
        )
    return f"""
<h2>Cronograma Gantt · semanas del proyecto</h2>
<p style="font-size: 10pt; color: #666;">Marca con X (o con un color en Excel) las semanas en las que se ejecuta cada actividad. Adapta el número de semanas a tu proyecto.</p>
<table class="matrix">
<tr><th style="width:4%;">#</th><th style="width:24%;">Actividad</th><th style="width:14%;">Responsable</th>{weeks_header}</tr>
{rows}
</table>
"""


def render_indicators_example(num):
    """Para 4.4.1 (reporte de evaluación de resultados), añade tabla antes/después."""
    if num != "4.4.1":
        return ""
    rows = ""
    for i in range(1, 7):
        rows += (
            f'<tr>\n'
            f'<td style="text-align:center;color:#28467e;font-weight:bold;">{i}</td>\n'
            f'<td class="fill">&nbsp;</td>\n'
            f'<td class="fill" style="text-align:center;">&nbsp;</td>\n'
            f'<td class="fill" style="text-align:center;">&nbsp;</td>\n'
            f'<td class="fill" style="text-align:center;color:#999;">=(actual-base)/base</td>\n'
            f'<td class="fill">&nbsp;</td>\n'
            f'</tr>\n'
        )
    return f"""
<h2>Tabla antes/después de indicadores</h2>
<p style="font-size: 10pt; color: #666;">Por cada indicador del proyecto, registra la línea base, el valor actual y la variación. Cumple con el primer criterio del F21 explícitamente.</p>
<table class="matrix">
<tr><th style="width:4%;">#</th><th style="width:30%;">Indicador</th><th style="width:14%;">Línea base</th><th style="width:14%;">Valor actual</th><th style="width:14%;">Variación %</th><th>Observación</th></tr>
{rows}
</table>
<h2 style="margin-top:24pt;">Registro de incidencias</h2>
<table class="matrix">
<tr><th style="width:4%;">#</th><th style="width:14%;">Fecha</th><th style="width:30%;">Incidencia</th><th style="width:24%;">Acción correctiva</th><th>Lección aprendida</th></tr>
{''.join(f'<tr><td style="text-align:center;color:#28467e;font-weight:bold;">{i}</td><td class="fill">&nbsp;</td><td class="fill">&nbsp;</td><td class="fill">&nbsp;</td><td class="fill">&nbsp;</td></tr>' for i in range(1, 5))}
</table>
"""




def render_product(p):
    """Devuelve (contenido, extensión). Para pptx, contenido es un objeto
    Presentation que main() guarda con .save(path)."""
    f21_html = render_f21(p["f21"])

    if p["format"] == "word":
        questions_html = render_questions_word(p["preguntas"])
        html = WORD_TEMPLATE.format(
            num=p["num"], title=p["title"], elemento=p["elemento"], intro=p["intro"],
            espiga=p["espiga"], f21_html=f21_html, questions_html=questions_html,
            base_styles=BASE_STYLES,
        )
        return html, "doc"
    elif p["format"] == "excel":
        questions_html = render_questions_excel(p["preguntas"])
        matrix_html = render_matrix_example(p["num"]) + render_gantt_example(p["num"]) + render_indicators_example(p["num"])
        html = EXCEL_TEMPLATE.format(
            num=p["num"], title=p["title"], elemento=p["elemento"], intro=p["intro"],
            espiga=p["espiga"], f21_html=f21_html, questions_html=questions_html,
            matrix_html=matrix_html, base_styles=BASE_STYLES,
        )
        return html, "xls"
    elif p["format"] == "pptx":
        return generate_pptx_capacitacion(p), "pptx"
    raise ValueError(f"Formato desconocido: {p['format']}")


def main():
    # Borrar archivos legacy: OOXML generados antes (.docx/.xlsx) que
    # ya no usamos, y carpetas auxiliares que Office crea al previsualizar
    # HTML-en-extensión-Office.
    for ext in (".docx", ".xlsx"):
        for f in OUT.glob(f"*{ext}"):
            f.unlink()
            print(f"  [del legacy] {f.name}")
    for d in OUT.glob("*.files"):
        if d.is_dir():
            import shutil
            shutil.rmtree(d)
            print(f"  [del legacy dir] {d.name}")
    # Para el 3.4.3: borrar versión .doc previa si existía (cambia a .pptx)
    old_343 = OUT / "3-4-3-material-capacitacion.doc"
    if old_343.exists():
        old_343.unlink()
        print(f"  [del legacy] {old_343.name}")

    for p in PRODUCTS:
        content, ext = render_product(p)
        filename = f"{p['num'].replace('.', '-')}-{p['slug']}.{ext}"
        filepath = OUT / filename
        if p["format"] == "pptx":
            content.save(filepath)
        else:
            filepath.write_text(content, encoding="utf-8")
        print(f"  [ok] {filename}  ({p['format'].upper()})")
    print(f"\nGenerados {len(PRODUCTS)} templates en {OUT}/")


if __name__ == "__main__":
    main()
