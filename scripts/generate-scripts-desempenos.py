#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Genera los PPTX de guion (script de video) por desempeno del EC
"Gestion de contenido de marketing digital con IA en la MiPyME".

Replica el formato del entregable de referencia
`extras/Script para desempeno 7_06082026.pptx`:

  1. Portada          -> titulo + duracion de video
  2. Notas de diseno  -> conteo de palabras, duracion, indicaciones de animacion
  3. Escena de apertura -> plano general + locucion de narrador + titulo en pantalla
  4..n Escenas        -> imagen del personaje + dialogo + direccion de actuacion
                         (expresion facial / postura corporal) + chip de criterio

El contenido vive en `scripts/guiones-desempenos.json` para que el equipo
editorial pueda ajustarlo sin tocar este archivo.

Uso:
    python scripts/generate-scripts-desempenos.py            # todos
    python scripts/generate-scripts-desempenos.py 1 6 8      # solo esos numeros

Salida: extras/scripts-desempenos/  (carpeta fuera de git)
"""

import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DATA = ROOT / "scripts" / "guiones-desempenos.json"
ASSETS = ROOT / "extras" / "script-assets"
OUTDIR = ROOT / "extras" / "scripts-desempenos"

# El guion de referencia se usa como PLANTILLA: de el heredamos tema (fuente
# Aptos y los accent del tema de Office), tamano de diapositiva y layouts.
# Asi los guiones nuevos salen visualmente hermanados con el original.
PLANTILLA = ROOT / "extras" / "Script para desempeño 7_06082026.pptx"

# Colores por rol, tal como los usa el guion de referencia:
#   dialogo y panel de direccion -> accent1   chip de criterio -> accent2
#   badge de numero de escena    -> accent6   texto sobre ellos -> blanco
ROL_DIALOGO = MSO_THEME_COLOR.ACCENT_1
ROL_DIRECCION = MSO_THEME_COLOR.ACCENT_1
ROL_CRITERIO = MSO_THEME_COLOR.ACCENT_2
ROL_BADGE = MSO_THEME_COLOR.ACCENT_6
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
TINTA = RGBColor(0x0E, 0x28, 0x41)  # dk2 del tema, para texto sobre fondo claro

SLIDE_W, SLIDE_H = 13.333, 7.5


def _layout(prs, nombre):
    """Layout por nombre, con respaldo por indice si la plantilla cambia de idioma."""
    for lay in prs.slide_layouts:
        if lay.name.lower().startswith(nombre.lower()):
            return lay
    return prs.slide_layouts[6]


def _blank(prs):
    return prs.slides.add_slide(_layout(prs, "En blanco"))


def _textbox(slide, left, top, width, height, *, fill=None, tema=None,
             shape=MSO_SHAPE.RECTANGLE, radius=None):
    """fill = RGBColor explicito; tema = MSO_THEME_COLOR (gana sobre fill)."""
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    if tema is not None:
        sh.fill.solid()
        sh.fill.fore_color.theme_color = tema
    elif fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    sh.line.fill.background()
    if radius is not None and sh.adjustments:
        sh.adjustments[0] = radius
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    return sh


def _write(tf, blocks, *, size=14, color=TINTA, align=PP_ALIGN.LEFT,
           anchor=MSO_ANCHOR.TOP, space_after=4):
    """blocks: lista de (texto, bold) o string suelto. La fuente se hereda del
    tema de la plantilla (Aptos) — nunca se fija font.name a proposito."""
    tf.vertical_anchor = anchor
    if isinstance(blocks, str):
        blocks = [(blocks, False)]
    first = True
    for item in blocks:
        text, bold = item if isinstance(item, (tuple, list)) else (item, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _linea(tf, tramos, *, size, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE):
    """Un solo parrafo con varios runs (para 'Criterio: <texto>' en una linea)."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for texto, bold in tramos:
        r = p.add_run()
        r.text = texto
        r.font.bold = bold
        r.font.size = Pt(size)
        r.font.color.rgb = color
    tf.vertical_anchor = anchor


def _badge(slide, numero):
    """Circulo con el numero de escena (accent6), esquina superior derecha.
    Se dibuja al final para quedar por encima del panel, como en el ejemplo."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.03), Inches(0.26),
                                Inches(0.87), Inches(0.81))
    sh.fill.solid()
    sh.fill.fore_color.theme_color = ROL_BADGE
    sh.line.fill.background()
    sh.shadow.inherit = False
    _write(sh.text_frame, [(str(numero), False)], size=32, color=BLANCO,
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _criterio(slide, texto):
    """Chip (accent2) con el criterio del estandar que cubre la escena."""
    sh = _textbox(slide, 0.77, 0.49, 3.26, 1.05, tema=ROL_CRITERIO,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    _linea(sh.text_frame, [("Criterio: ", True), (texto, False)], size=12)


def _dialogo(slide, hablante, texto, *, narrador=False):
    """Barra de dialogo inferior (accent1). La imagen del personaje queda
    detras, como en el guion de referencia."""
    sh = _textbox(slide, 0.44, 5.45, 12.45, 1.62, tema=ROL_DIALOGO)
    etiqueta = "En audio (narrador):" if narrador else f"En audio (voz {hablante}):"
    _write(sh.text_frame, [(etiqueta, True), (texto, False)],
           size=14, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE, space_after=3)


def _direccion(slide, facial, postura):
    """Panel de direccion de actuacion (accent1), columna derecha.
    El margen superior deja libre la esquina que ocupa el badge de escena."""
    sh = _textbox(slide, 10.04, 0.54, 3.10, 4.91, tema=ROL_DIRECCION)
    sh.text_frame.margin_top = Inches(0.62)
    _write(sh.text_frame,
           [("Expresión facial:", True), (facial, False),
            ("", False),
            ("Postura corporal:", True), (postura, False)],
           size=11.5, color=BLANCO, space_after=3)


def _portada(prs, g):
    """Usa el layout 'Diapositiva de titulo' de la plantilla, igual que el ejemplo."""
    s = prs.slides.add_slide(_layout(prs, "Diapositiva de t"))
    titulo = sub = None
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            titulo = ph
        elif sub is None:
            sub = ph
    if titulo is not None:
        titulo.text_frame.text = f"Script para desempeño {g['numero']}"
    if sub is not None:
        tf = sub.text_frame
        tf.text = f"Duración de video: {g['duracion']}"
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = g["titulo"]
        r.font.size = Pt(18)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = g["elemento"]
        r2.font.size = Pt(14)
    return s


def _notas(prs, g, palabras):
    """Usa el layout 'Titulo y objetos', igual que el ejemplo."""
    s = prs.slides.add_slide(_layout(prs, "Título y objetos"))
    cuerpo_ph = None
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.text = "Notas diseño"
        elif cuerpo_ph is None:
            cuerpo_ph = ph

    cuerpo = [
        (f"Palabras: {palabras}", False),
        (f"Duración de narración: {g['duracion']}", False),
        ("Cambiar las expresiones y postura de los personajes para dar movimiento al video.", False),
        ("", False),
        ("Criterios del estándar cubiertos en este guion:", True),
    ]
    for c in g["criterios"]:
        cuerpo.append((f"•  {c}", False))
    if g.get("ahv"):
        cuerpo.append(("", False))
        cuerpo.append(("AHV que se observan:", True))
        for a in g["ahv"]:
            cuerpo.append((f"•  {a}", False))
    if g.get("nota_arte"):
        cuerpo.append(("", False))
        cuerpo.append(("Nota de arte:", True))
        cuerpo.append((g["nota_arte"], False))

    if cuerpo_ph is not None:
        _write(cuerpo_ph.text_frame, cuerpo, size=13, color=TINTA, space_after=3)
    else:
        box = _textbox(s, 0.7, 1.6, 11.9, 5.3)
        _write(box.text_frame, cuerpo, size=13, color=TINTA, space_after=3)
    return s


def _apertura(prs, g, esc, n):
    """Plano general con el rotulo de titulo en pantalla, como el ejemplo."""
    s = _blank(prs)
    img = ASSETS / esc.get("imagen", "escena-wide.png")
    if img.exists():
        ancho = 11.79
        with Image.open(img) as im:
            alto = ancho * im.height / im.width
        s.shapes.add_picture(str(img), Inches(0.68), Inches(0.94),
                             width=Inches(ancho), height=Inches(alto))

    titulo = _textbox(s, 4.12, 0.28, 5.05, 0.58, fill=BLANCO)
    _linea(titulo.text_frame,
           [("Título en pantalla: ", True), (esc["titulo_pantalla"], False)],
           size=15, color=TINTA)
    titulo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    _dialogo(s, None, esc["narrador"], narrador=True)
    _badge(s, n)
    return s


def _escena(prs, esc, n):
    """Chip de criterio a la izquierda, personaje al centro, direccion a la derecha.
    El personaje se dibuja primero para que la barra de dialogo lo tape por
    abajo, que es el efecto del guion de referencia."""
    s = _blank(prs)
    hablante = esc["hablante"]
    img_name = esc.get("imagen") or ("martin.png" if hablante.lower().startswith("mart")
                                     else "mariana.png")
    img = ASSETS / img_name

    if img.exists():
        alto = 6.00
        with Image.open(img) as im:
            ancho = alto * im.width / im.height
        # centrado en la banda libre entre el chip y el panel de direccion
        left = 4.15 + (5.85 - ancho) / 2
        s.shapes.add_picture(str(img), Inches(left), Inches(0.49),
                             width=Inches(ancho), height=Inches(alto))

    _direccion(s, esc["facial"], esc["postura"])
    if esc.get("criterio"):
        _criterio(s, esc["criterio"])
    _dialogo(s, hablante, esc["dialogo"])
    _badge(s, n)
    return s


def contar_palabras(g):
    total = 0
    for esc in g["escenas"]:
        total += len((esc.get("narrador") or esc.get("dialogo") or "").split())
    return total


def _nueva_presentacion():
    """Abre el guion de referencia como plantilla y le quita las diapositivas,
    conservando tema, fuentes, layouts y tamano. Si no esta, cae al default."""
    if not PLANTILLA.exists():
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        return prs, False

    prs = Presentation(str(PLANTILLA))
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        prs.part.drop_rel(sld.rId)
        lst.remove(sld)
    return prs, True


def construir(g):
    prs, _ = _nueva_presentacion()

    palabras = contar_palabras(g)
    _portada(prs, g)
    _notas(prs, g, palabras)

    for i, esc in enumerate(g["escenas"], start=1):
        if esc.get("tipo") == "apertura":
            _apertura(prs, g, esc, i)
        else:
            _escena(prs, esc, i)

    OUTDIR.mkdir(parents=True, exist_ok=True)
    slug = g["titulo"].lower()
    for a, b in (("á", "a"), ("é", "e"), ("í", "i"), ("ó", "o"), ("ú", "u"),
                 ("ñ", "n"), (" ", "-"), (",", ""), (":", ""), ("/", "-")):
        slug = slug.replace(a, b)
    slug = "-".join(filter(None, slug.split("-")))[:60]
    destino = OUTDIR / f"Script para desempeno {g['numero']} - {slug}.pptx"
    prs.save(str(destino))
    return destino, palabras, len(g["escenas"])


def main():
    if not DATA.exists():
        sys.exit(f"No encuentro {DATA}")
    guiones = json.loads(DATA.read_text(encoding="utf-8"))

    filtro = {int(a) for a in sys.argv[1:] if a.isdigit()}
    if filtro:
        guiones = [g for g in guiones if g["numero"] in filtro]

    if not ASSETS.exists():
        print(f"Aviso: no existe {ASSETS} — los PPTX saldran sin imagenes.\n")
    if not PLANTILLA.exists():
        print(f"Aviso: no encuentro la plantilla {PLANTILLA.name};\n"
              f"       los guiones saldran con el tema default de PowerPoint.\n")

    for g in guiones:
        destino, palabras, escenas = construir(g)
        print(f"[{g['numero']}] {destino.name}")
        print(f"     {escenas} escenas · {palabras} palabras · {g['duracion']}")


if __name__ == "__main__":
    main()
